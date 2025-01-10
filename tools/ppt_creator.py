import os
import csv
from pptx import Presentation
from pptx.util import Inches

def create_presentation(base_folder, output_filename):
    # Initialize PowerPoint presentation
    presentation = Presentation()

    # Collect folders with seed numbers
    folders = []
    for folder_name in os.listdir(base_folder):
        if folder_name.startswith("results_"):
            # Extract the seed number from the folder name
            seed_number = int(folder_name.split("_")[2])  # Convert to integer for sorting
            folder_path = os.path.join(base_folder, folder_name)
            folders.append((seed_number, folder_path))

    # Sort folders by seed number in ascending order
    folders.sort(key=lambda x: x[0])

    # Iterate over sorted folders
    for seed_number, folder_path in folders:
        # Paths for required images and CSV file
        it_0_path = os.path.join(folder_path, "it_0.png")
        it_1000_path = os.path.join(folder_path, "it_1000.png")
        aesthetic_evolution_path = os.path.join(folder_path, "score_evolution.png")
        loss_evolution_path = os.path.join(folder_path, "loss_evolution.png")
        csv_path = os.path.join(folder_path, "score_results.csv")

        # Extract scores from CSV
        score_initial = None
        score_final = None
        if os.path.exists(csv_path):
            with open(csv_path, 'r') as csvfile:
                reader = csv.DictReader(csvfile)
                for row in reader:
                    iteration = int(row['iteration'])
                    score = float(row['score'])
                    if iteration == 0:
                        score_initial = score
                    elif iteration == 1000:
                        score_final = score

        # Slide 1: it_0.png and it_1000.png
        if os.path.exists(it_0_path) and os.path.exists(it_1000_path):
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide
            title = slide.shapes.title
            title.text = f"Seed {seed_number}"

            # Add it_0.png
            slide.shapes.add_picture(it_0_path, Inches(0.5), Inches(2), height=Inches(4))

            # Add legend below it_0.png
            left = Inches(0.5)
            top = Inches(6.2)
            width = Inches(4)
            textbox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
            text = "Initial iteration"
            if score_initial is not None:
                text += f"\nScore: {score_initial:.4f}"
            textbox.text = text

            # Add it_1000.png
            slide.shapes.add_picture(it_1000_path, Inches(5.5), Inches(2), height=Inches(4))

            # Add legend below it_1000.png
            left = Inches(5.5)
            top = Inches(6.2)
            width = Inches(4)
            textbox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
            text = "Final iteration"
            if score_final is not None:
                text += f"\nScore: {score_final:.4f}"
            textbox.text = text

        # Slide 2: aesthetic_evolution.png and loss_evolution.png
        if os.path.exists(aesthetic_evolution_path) and os.path.exists(loss_evolution_path):
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide
            title = slide.shapes.title
            title.text = f"Seed {seed_number}"

            # Add aesthetic_evolution.png
            slide.shapes.add_picture(aesthetic_evolution_path, Inches(0), Inches(2), height=Inches(4))

            # Add loss_evolution.png
            slide.shapes.add_picture(loss_evolution_path, Inches(5), Inches(2), height=Inches(4))

    # Save the presentation
    presentation.save(output_filename)
    print(f"Presentation saved as {output_filename}")

# Example usage
base_folder = "results/adam_embedding_laion"  # Replace with your base folder path
output_filename = "results/adam_embedding_laion/summary.pptx"  # Replace with your desired output file path
create_presentation(base_folder, output_filename)
